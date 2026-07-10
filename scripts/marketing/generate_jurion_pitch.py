#!/usr/bin/env python3
"""Genera assets visuales, pitch deck Jurion (.pptx) y one-pager (.pdf)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "docs" / "pitch"
ASSETS_DIR = OUT_DIR / "assets"
HTML_ONEPAGER = OUT_DIR / "JurionOnePager.html"
HTML_VALOR = OUT_DIR / "cuadro-valor-ahorro-completo.html"
ASSET_TEMPLATES = OUT_DIR / "asset-templates.html"

BG = RGBColor(0x05, 0x0D, 0x1A)
BG_SOFT = RGBColor(0x0A, 0x18, 0x2E)
BLUE_INST = RGBColor(0x04, 0x2C, 0x53)
BLUE_ACCENT = RGBColor(0x37, 0x8A, 0xDD)
TEXT = RGBColor(0xF0, 0xF4, 0xFF)
TEXT_MUTED = RGBColor(0x88, 0x99, 0xBB)
SUCCESS = RGBColor(0x22, 0xC5, 0x5E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.65)

ASSET_IDS = ("hero", "dashboard", "correo", "trabajo", "flow", "ia", "oficios")


def generate_assets() -> dict[str, Path]:
    if not ASSET_TEMPLATES.exists():
        raise FileNotFoundError(ASSET_TEMPLATES)

    from playwright.sync_api import sync_playwright

    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    uri = ASSET_TEMPLATES.resolve().as_uri()
    paths: dict[str, Path] = {}

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1400, "height": 900})
        page.goto(uri, wait_until="networkidle")
        for asset_id in ASSET_IDS:
            out = ASSETS_DIR / f"{asset_id}.png"
            page.locator(f"#{asset_id}").screenshot(path=str(out))
            paths[asset_id] = out
        browser.close()

    return paths


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _shape(slide, kind, left, top, width, height, fill: RGBColor, line: RGBColor | None = None, transparency: float = 0):
    s = slide.shapes.add_shape(kind, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if transparency:
        s.fill.transparency = transparency
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def _rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None, transparency: float = 0):
    return _shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height, fill, line, transparency)


def _oval(slide, left, top, width, height, fill: RGBColor, transparency: float = 0.35):
    return _shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left, top, width, height, fill, transparency=transparency)


def _accent_top(slide) -> None:
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.055), BLUE_ACCENT)


def _glows(slide) -> None:
    _oval(slide, Inches(9.5), Inches(-0.8), Inches(4.5), Inches(4.5), BLUE_ACCENT, 0.55)
    _oval(slide, Inches(-1.2), Inches(4.8), Inches(3.8), Inches(3.8), BLUE_INST, 0.45)


def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _p(tf, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=6, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.name = "Calibri Light"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def _label(slide, text: str, left=MARGIN, top=Inches(0.48)):
    box = _tb(slide, left, top, Inches(5), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    _p(tf, text.upper(), size=10, bold=True, color=BLUE_ACCENT, space_after=0, first=True)
    return box


def _title(slide, text: str, left=MARGIN, top=Inches(0.82), width=None):
    w = width or (SLIDE_W - MARGIN * 2)
    box = _tb(slide, left, top, w, Inches(1.35))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    _p(tf, text, size=36, bold=True, color=TEXT, space_after=0, first=True)
    return box


def _subtitle(slide, text: str, left=MARGIN, top=Inches(2.05), width=None):
    w = width or (SLIDE_W - MARGIN * 2)
    box = _tb(slide, left, top, w, Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    _p(tf, text, size=17, color=TEXT_MUTED, space_after=0, first=True)
    return box


def _footer(slide, n: int, total: int) -> None:
    _rect(slide, Inches(0), SLIDE_H - Inches(0.42), SLIDE_W, Inches(0.015), BLUE_ACCENT)
    box = _tb(slide, MARGIN, SLIDE_H - Inches(0.38), SLIDE_W - MARGIN * 2, Inches(0.28))
    tf = box.text_frame
    tf.clear()
    _p(tf, "Jurion", size=10, bold=True, color=TEXT_MUTED, space_after=0, first=True)
    p2 = tf.add_paragraph()
    p2.text = f"{n} / {total}"
    p2.font.name = "Calibri Light"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.RIGHT


def _add_image(slide, path: Path, left, top, width, height=None):
    if not path.exists():
        ph = _rect(slide, left, top, width, height or Inches(3.5), BG_SOFT, BLUE_ACCENT)
        tb = _tb(slide, left + Inches(0.2), top + Inches(1.2), width - Inches(0.4), Inches(1))
        tf = tb.text_frame
        tf.clear()
        _p(tf, "Vista Jurion", size=16, bold=True, color=BLUE_ACCENT, align=PP_ALIGN.CENTER, first=True)
        return ph
    _rect(slide, left - Inches(0.03), top - Inches(0.03), width + Inches(0.06), Inches(4.2), BG_SOFT, BLUE_ACCENT)
    pic = slide.shapes.add_picture(str(path), left, top, width=width)
    if height:
        pic.height = height
    return pic


def _card(slide, left, top, w, h, title, body, featured=False):
    fill = BLUE_INST if featured else BG_SOFT
    line = BLUE_ACCENT
    _rect(slide, left, top, w, h, fill, line, transparency=0.08 if featured else 0)
    box = _tb(slide, left + Inches(0.18), top + Inches(0.15), w - Inches(0.36), h - Inches(0.25))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    _p(tf, title, size=14, bold=True, color=TEXT if not featured else BLUE_ACCENT, space_after=6, first=True)
    _p(tf, body, size=11, color=TEXT_MUTED, space_after=0)


def build_pptx(path: Path, assets: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    total = 14

    # 1 Cover — full visual
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _add_image(s, assets["hero"], Inches(6.8), Inches(0.55), Inches(6.2))
    box = _tb(s, MARGIN, Inches(1.6), Inches(6.0), Inches(1.0))
    tf = box.text_frame
    tf.clear()
    _p(tf, "Jurion.", size=54, bold=True, color=TEXT, first=True)
    _title(s, "La Rama estaba en mil herramientas.", top=Inches(2.55), width=Inches(6.2))
    sub = _tb(s, MARGIN, Inches(3.55), Inches(6.0), Inches(1.5))
    tf = sub.text_frame
    tf.word_wrap = True
    tf.clear()
    _p(tf, "Ahora, un solo despacho digital.", size=28, bold=True, color=BLUE_ACCENT, first=True)
    _p(tf, "Outlook · SGDE · plazos · expediente · IA judicial", size=14, color=TEXT_MUTED, space_after=0)
    badge = _rect(s, MARGIN, Inches(5.5), Inches(5.8), Inches(0.42), BG_SOFT, BLUE_ACCENT)
    btb = _tb(s, MARGIN + Inches(0.15), Inches(5.58), Inches(5.5), Inches(0.3))
    tf = btb.text_frame
    tf.clear()
    _p(tf, "PCSJA24-12243  ·  Ley 2213/2022  ·  Colombia 2026", size=11, color=BLUE_ACCENT, first=True)
    _footer(s, 1, total)

    # 2 Problem — split
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "El problema")
    _title(s, "¿Cuántas ventanas abre un juzgado\npara un solo proceso?", width=Inches(5.8))
    _subtitle(s, "Correo, SGDE, Word, Excel y términos en sistemas distintos.", width=Inches(5.5))
    cards = [
        ("Vencimientos invisibles", "Un día de más puede significar nulidad o revictimización."),
        ("Correo sin procesar", "Repartos y memorials que alguien debe clasificar a mano."),
        ("SGDE desconectado", "Alfresco vs Word vs Outlook — nadie sincroniza."),
    ]
    for i, (t, b) in enumerate(cards):
        _card(s, MARGIN, Inches(3.35 + i * 1.22), Inches(5.6), Inches(1.05), t, b)
    _add_image(s, assets["dashboard"], Inches(6.85), Inches(1.0), Inches(6.0))
    _footer(s, 2, total)

    # 3 Unification
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "La solución")
    _title(s, "Jurion concentra todo el flujo")
    _add_image(s, assets["flow"], MARGIN, Inches(2.15), SLIDE_W - MARGIN * 2)
    _footer(s, 3, total)

    # 4 Dashboard full bleed
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _accent_top(s)
    _label(s, "El producto real")
    _title(s, "Resumen operativo del despacho", top=Inches(0.78))
    _add_image(s, assets["dashboard"], MARGIN, Inches(2.0), SLIDE_W - MARGIN * 2)
    cap = _tb(s, MARGIN, Inches(6.55), SLIDE_W - MARGIN * 2, Inches(0.35))
    tf = cap.text_frame
    tf.clear()
    _p(tf, "Semáforo de términos · pendientes de ingreso · actividad reciente en un solo tablero", size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER, first=True)
    _footer(s, 4, total)

    # 5 Correo + IA
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Correo judicial")
    _title(s, "Recibir, enviar y radicar\ndesde Outlook M365", width=Inches(5.5))
    _subtitle(s, "Bandeja integrada: analizar con IA, aprobar ingreso y enviar oficios a las partes.", width=Inches(5.3))
    _card(s, MARGIN, Inches(3.5), Inches(5.4), Inches(1.1), "Correo entrante", "Reparto, memorial, respuesta — clasificación IA y radicación .eml.", featured=True)
    _card(s, MARGIN, Inches(4.75), Inches(5.4), Inches(1.1), "Correo saliente", "Envío de oficios y comunicaciones a partes sin salir de Jurion.", featured=False)
    _add_image(s, assets["correo"], Inches(6.55), Inches(0.85), Inches(6.35))
    _footer(s, 5, total)

    # 6 De punta a punta — oficios
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Flujo completo")
    _title(s, "Ingresa el reparto, sale el oficio", top=Inches(0.78))
    _subtitle(
        s,
        "Plantillas con datos del expediente · envío Outlook M365 · registro de notificación · trazabilidad.",
        top=Inches(1.65),
    )
    _add_image(s, assets["oficios"], MARGIN, Inches(2.05), SLIDE_W - MARGIN * 2)
    steps = [
        "Correo reparto → IA clasifica → humano aprueba",
        "Expediente + documentos desde plantillas CSJ",
        "Auto o fallo firmado → tarea a secretaría",
        "Oficio generado → enviado por correo → registrado",
    ]
    gap = (SLIDE_W - MARGIN * 2 - Inches(0.36)) / 4
    for i, step in enumerate(steps):
        left = MARGIN + i * (gap + Inches(0.12))
        _rect(s, left, Inches(6.35), gap, Inches(0.55), BG_SOFT, SUCCESS if i == 3 else BLUE_ACCENT)
        tb = _tb(s, left + Inches(0.1), Inches(6.45), gap - Inches(0.2), Inches(0.38))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        _p(tf, step, size=9, bold=(i == 3), color=TEXT if i < 3 else SUCCESS, align=PP_ALIGN.CENTER, first=True)
    _footer(s, 6, total)

    # 7 Centro trabajo
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _accent_top(s)
    _add_image(s, assets["trabajo"], MARGIN, Inches(0.55), Inches(7.8))
    _label(s, "Centro de trabajo", left=Inches(8.85))
    _title(s, "Tareas automáticas\npor rol", left=Inches(8.85), top=Inches(0.82), width=Inches(4.0))
    _subtitle(s, "Informe, auto, fallo, SGDE y revisión — asignados al sustanciador correcto.", left=Inches(8.85), top=Inches(2.35), width=Inches(3.9))
    items = ["Notificaciones en tiempo real", "Motor de etapas procesales", "Revisión juez con trazabilidad"]
    box = _tb(s, Inches(8.85), Inches(3.55), Inches(3.9), Inches(2.5))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        _p(tf, f"✦  {item}", size=15, color=TEXT if i == 0 else TEXT_MUTED, space_after=10, first=(i == 0))
    _footer(s, 7, total)

    # 8 IA jewel
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _add_image(s, assets["ia"], MARGIN, Inches(0.7), SLIDE_W - MARGIN * 2)
    _footer(s, 8, total)

    # 9 Features grid
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Bondades en producción")
    _title(s, "Todo lo que el despacho necesita hoy", top=Inches(0.78))
    feats = [
        ("Oficios y notificaciones", "Plantillas + tareas a secretaría tras firma del auto o fallo."),
        ("Correo entrante y saliente", "Outlook M365: recibir, enviar oficios y radicar .eml."),
        ("SGDE nativo", "Crear, vincular y sincronizar piezas bidireccional."),
        ("Documentos CSJ", "Informe, auto admisorio y Word desde plantillas del despacho."),
        ("Semáforo términos", "EN TÉRMINO · URGENTE · VENCIDO."),
        ("Auditoría", "Actuaciones, envíos y cambios por usuario y hora."),
    ]
    gap = Inches(0.18)
    cw = (SLIDE_W - MARGIN * 2 - gap * 2) / 3
    ch = Inches(1.35)
    for i, (t, b) in enumerate(feats):
        r, c = divmod(i, 3)
        _card(s, MARGIN + c * (cw + gap), Inches(2.15) + r * (ch + gap), cw, ch, t, b, featured=(i == 0))
    _footer(s, 9, total)

    # 10 Diferenciadores — texto + imagen lateral
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Por qué Jurion")
    _title(s, "Hecho para la Rama,\npor quienes la conocen", width=Inches(5.8))
    diffs = [
        ("01", "Operadores judiciales", "Diseñado con jueces y equipos en activo."),
        ("02", "Norma nativa", "PCSJA24-12243 y Ley 2213 — no es un add-on."),
        ("03", "Desde la trinchera", "Cada pantalla responde a un dolor real."),
        ("04", "Memoria propia", "Precedentes indexados gobernados por el juez."),
    ]
    for i, (num, t, b) in enumerate(diffs):
        r, c = divmod(i, 2)
        left = MARGIN + c * Inches(3.05)
        top = Inches(2.45) + r * Inches(1.45)
        _rect(s, left, top, Inches(2.85), Inches(1.25), BG_SOFT, BLUE_ACCENT)
        tb = _tb(s, left + Inches(0.15), top + Inches(0.12), Inches(2.55), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        _p(tf, num, size=18, bold=True, color=BLUE_ACCENT, first=True)
        _p(tf, t, size=12, bold=True, color=TEXT, space_after=3)
        _p(tf, b, size=10, color=TEXT_MUTED, space_after=0)
    _add_image(s, assets["dashboard"], Inches(6.55), Inches(0.75), Inches(6.35))
    _footer(s, 10, total)

    # 11 Compliance
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Alineación institucional")
    _title(s, "Coherencia con el protocolo\nde la Rama Judicial", width=Inches(6.0))
    panel = _rect(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.8), BLUE_INST, BLUE_ACCENT, 0.05)
    ptb = _tb(s, Inches(7.55), Inches(2.0), Inches(4.8), Inches(3.8))
    tf = ptb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate([
        "PDF conforme CENDOJ",
        "Nombres TitleCase sin guiones",
        "Cuadernos C01, C02 alineados SGDE",
        "InformeIngresoDespacho.pdf",
        "Ley 2213 · Decreto 2591",
    ]):
        _p(tf, line, size=16, bold=(i == 0), color=TEXT if i < 4 else BLUE_ACCENT, space_after=12, first=(i == 0))
    bullets = _tb(s, MARGIN, Inches(2.4), Inches(6.2), Inches(3.5))
    tf = bullets.text_frame
    tf.clear()
    _p(tf, "No es otra app más.", size=22, bold=True, color=TEXT, first=True)
    _p(tf, "Implementa las mismas reglas que la Rama promueve para interoperabilidad, conservación y expediente digital.", size=15, color=TEXT_MUTED, space_after=0)
    _footer(s, 11, total)

    # 12 Roadmap
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Roadmap")
    _title(s, "Transformación digital 2026–2028")
    roadmap = [
        ("Q2 2026", "Beta tutelas", "Correo + IA + SGDE"),
        ("Q3 2026", "Civil", "Citación asistida"),
        ("Q4 2026", "Todas las jurisdicciones", "Penal · laboral · familia"),
        ("2027", "Expansión nacional", "Costa a costa"),
        ("2028", "Rama unificada", "Una arquitectura"),
    ]
    _rect(s, MARGIN, Inches(4.2), SLIDE_W - MARGIN * 2, Inches(0.04), BLUE_ACCENT)
    gap = (SLIDE_W - MARGIN * 2) / 5
    for i, (when, title, body) in enumerate(roadmap):
        cx = MARGIN + gap * i + gap / 2
        col = SUCCESS if i == 0 else BLUE_ACCENT
        _oval(s, cx - Inches(0.14), Inches(4.05), Inches(0.28), Inches(0.28), col, 0)
        tb = _tb(s, cx - Inches(1.0), Inches(2.4 if i % 2 == 0 else 4.55), Inches(2.0), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        _p(tf, when, size=12, bold=True, color=col, align=PP_ALIGN.CENTER, first=True)
        _p(tf, title, size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER, space_after=4)
        _p(tf, body, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER, space_after=0)
    _footer(s, 12, total)

    # 13 Ahorro de tiempo y valor institucional
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _label(s, "Impacto operativo")
    _title(s, "Lo más importante para el despacho", top=Inches(0.78))
    _subtitle(
        s,
        "Gobernanza y estadística primero; después IA (precedentes, resúmenes) y términos automáticos.",
        top=Inches(1.72),
    )
    primary = [
        (
            "Trazabilidad y auditoría",
            "Actuaciones + historial: quién cambió qué, cuándo y en qué pieza. Revisiones Word registradas.",
            "Rendir cuentas con registro",
        ),
        (
            "Estadística SIERJU automática",
            "Movimiento de Tutelas por periodo: derecho tutelado, ingresos y decisiones sin rellenar Excel a mano.",
            "Cierre mensual en minutos",
        ),
    ]
    pw = (SLIDE_W - MARGIN * 2 - Inches(0.14)) / 2
    for i, (t, b, save) in enumerate(primary):
        left = MARGIN + i * (pw + Inches(0.14))
        _rect(s, left, Inches(2.2), pw, Inches(1.35), BLUE_INST, SUCCESS, 0.05)
        tb = _tb(s, left + Inches(0.15), Inches(2.32), pw - Inches(0.3), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        _p(tf, t, size=14, bold=True, color=SUCCESS, first=True)
        _p(tf, b, size=10, color=TEXT_MUTED, space_after=5)
        _p(tf, save, size=11, bold=True, color=BLUE_ACCENT, space_after=0)
    secondary = [
        ("Precedentes + similitud", "~30–90 min → 3–10 min"),
        ("Lectura rápida IA / doc.", "~15–45 min → 2–5 min"),
        ("Oficios + envío correo", "Plantilla → Outlook → registro"),
        ("Términos automáticos", "Vigilancia 24/7"),
    ]
    sw = (SLIDE_W - MARGIN * 2 - Inches(0.36)) / 4
    for i, (t, save) in enumerate(secondary):
        left = MARGIN + i * (sw + Inches(0.12))
        _rect(s, left, Inches(3.72), sw, Inches(0.95), BG_SOFT, BLUE_ACCENT)
        tb = _tb(s, left + Inches(0.1), Inches(3.82), sw - Inches(0.2), Inches(0.75))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        _p(tf, t, size=11, bold=True, color=TEXT, first=True)
        _p(tf, save, size=10, bold=True, color=SUCCESS, space_after=0)
    rows = [
        ("Reparto nuevo", "~10 min", "~3–5 min", "~5–7 min"),
        ("Memorial / pieza", "~10 min", "~2–4 min", "~6–8 min"),
        ("Oficios + envío", "Word + Outlook aparte", "En Jurion", "Plantilla → correo → registro"),
        ("Sync SGDE", "~10–15 min", "~4–6 min", "~5–9 min"),
        ("Cierre SIERJU", "Horas en Excel", "Automático", "Desde expediente"),
    ]
    table_top = Inches(4.85)
    col_w = [Inches(3.2), Inches(1.5), Inches(1.5), Inches(3.8)]
    headers = ["Actividad", "Sin Jurion", "Con Jurion", "Ahorro / valor"]
    hx = MARGIN
    for hdr, w in zip(headers, col_w):
        _rect(s, hx, table_top, w - Inches(0.04), Inches(0.34), BG_SOFT, BLUE_ACCENT)
        hb = _tb(s, hx + Inches(0.08), table_top + Inches(0.05), w - Inches(0.16), Inches(0.26))
        tf = hb.text_frame
        tf.clear()
        _p(tf, hdr, size=9, bold=True, color=BLUE_ACCENT, first=True)
        hx += w
    for ri, row in enumerate(rows):
        y = table_top + Inches(0.36) + ri * Inches(0.42)
        x = MARGIN
        for ci, (cell, w) in enumerate(zip(row, col_w)):
            _rect(s, x, y, w - Inches(0.04), Inches(0.38), BG_SOFT if ri % 2 == 0 else BG, BLUE_ACCENT if ci == 0 else None)
            cb = _tb(s, x + Inches(0.08), y + Inches(0.06), w - Inches(0.16), Inches(0.28))
            tf = cb.text_frame
            tf.clear()
            _p(
                tf,
                cell,
                size=10 if ci == 0 else 9,
                bold=(ci == 0),
                color=SUCCESS if ri >= 3 and ci == 3 else TEXT if ci < 3 else TEXT_MUTED,
                first=True,
            )
            x += w
    note = _tb(s, MARGIN, Inches(7.05), SLIDE_W - MARGIN * 2, Inches(0.28))
    tf = note.text_frame
    tf.clear()
    _p(
        tf,
        "Estimación pre-piloto · trazabilidad y SIERJU son valor institucional; cifras de tiempo sujetas a medición",
        size=10,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        first=True,
    )
    _footer(s, 13, total)

    # 14 CTA — no contact
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s, BG)
    _glows(s)
    _accent_top(s)
    _add_image(s, assets["hero"], MARGIN, Inches(0.65), Inches(7.5))
    panel = _rect(s, Inches(8.45), Inches(1.2), Inches(4.25), Inches(5.2), BLUE_INST, BLUE_ACCENT, 0.05)
    tb = _tb(s, Inches(8.75), Inches(2.0), Inches(3.65), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    _p(tf, "Actualicemos la justicia juntos", size=28, bold=True, color=TEXT, first=True)
    _p(tf, "La Rama merece herramientas que hablen el mismo idioma.", size=14, color=TEXT_MUTED, space_after=12)
    _p(tf, "Correo · expediente · plazos · SGDE · IA judicial", size=13, bold=True, color=BLUE_ACCENT, space_after=12)
    _p(tf, "Hecho con quienes viven el despacho cada día.", size=13, color=TEXT_MUTED, space_after=0)
    _footer(s, 14, total)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def build_pdf_from_html(html_path: Path, pdf_path: Path) -> None:
    from playwright.sync_api import sync_playwright

    uri = html_path.resolve().as_uri()
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 816, "height": 1056})
        page.goto(uri, wait_until="networkidle")
        page.emulate_media(media="print")
        page.pdf(
            path=str(pdf_path),
            format="Letter",
            print_background=True,
            prefer_css_page_size=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
        browser.close()


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    print("Generando assets visuales…")
    assets = generate_assets()
    for k, v in assets.items():
        print(f"  {k}: {v}")

    pptx_path = OUT_DIR / "JurionPitch.pptx"
    pdf_path = OUT_DIR / "JurionOnePager.pdf"

    print("Generando pitch deck…")
    try:
        build_pptx(pptx_path, assets)
    except PermissionError:
        pptx_path = OUT_DIR / "JurionPitch-v2.pptx"
        print(f"  (archivo anterior abierto; guardando en {pptx_path.name})")
        build_pptx(pptx_path, assets)

    print("Generando PDF one-pager…")
    build_pdf_from_html(HTML_ONEPAGER, pdf_path)

    valor_pdf = OUT_DIR / "JurionCuadroValorAhorro.pdf"
    print("Generando PDF cuadro valor y ahorro…")
    build_pdf_from_html(HTML_VALOR, valor_pdf)

    print(f"\nListo:\n  {pptx_path}\n  {HTML_ONEPAGER}\n  {pdf_path}\n  {HTML_VALOR}\n  {valor_pdf}")


if __name__ == "__main__":
    main()
